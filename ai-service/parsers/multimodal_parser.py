import os
import pytesseract
from PIL import Image
import docx
import openpyxl
from pptx import Presentation
import pdfplumber
from parsers.base import FileParser
from providers.gemini_provider import GeminiProvider

class MultimodalParser(FileParser):
    def __init__(self):
        super().__init__()
        # Try to find tesseract on common windows paths if not in PATH
        if os.name == 'nt':
            common_paths = [
                r'C:\Program Files\Tesseract-OCR\tesseract.exe',
                r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
                r'C:\Users\sneha\AppData\Local\Tesseract-OCR\tesseract.exe'
            ]
            for p in common_paths:
                if os.path.exists(p):
                    pytesseract.pytesseract.tesseract_cmd = p
                    break

    def parse(self, file_path: str) -> str:
        ext = file_path.lower().split('.')[-1]
        
        try:
            if ext in ['png', 'jpg', 'jpeg', 'bmp', 'tiff']:
                return self._parse_image(file_path)
            elif ext == 'pdf':
                return self._parse_pdf(file_path)
            elif ext == 'docx':
                return self._parse_docx(file_path)
            elif ext == 'pptx':
                return self._parse_pptx(file_path)
            elif ext in ['xlsx', 'xls']:
                return self._parse_excel(file_path)
            elif ext in ['csv', 'tsv']:
                return self._parse_delimited(file_path, '\t' if ext == 'tsv' else ',')
            elif ext in ['txt', 'md', 'ts', 'tsx', 'js', 'jsx', 'json', 'yml', 'yaml', 'html', 'css', 'c', 'cpp', 'java', 'xml']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            else:
                raise ValueError(f"Unsupported file type: {ext}")
        except Exception as e:
            raise RuntimeError(f"Unable to parse {os.path.basename(file_path)}: {e}") from e

    def _parse_image(self, file_path: str) -> str:
        text = ""
        try:
            img = Image.open(file_path)
            ocr_text = pytesseract.image_to_string(img).strip()
            if ocr_text:
                text += f"OCR Text:\n{ocr_text}\n\n"
        except Exception as e:
            text += f"[OCR Error: {str(e)}]\n"
            
        try:
            gemini = GeminiProvider()
            if gemini.is_available():
                import google.generativeai as genai
                from PIL import Image as PILImage
                
                # Setup gemini vision call natively if possible, or just upload
                # Actually, the python google-generativeai supports passing PIL images directly
                model = gemini.client.aio.models if hasattr(gemini.client, 'aio') else None
                if model:
                    import asyncio
                    async def get_desc():
                        pil_img = PILImage.open(file_path)
                        response = await model.generate_content(
                            model="gemini-2.5-flash", 
                            contents=["Describe this image in extreme detail, including objects, colors, text, people, setting, clothing, and overall scene. List visually descriptive keywords.", pil_img]
                        )
                        return response.text
                    
                    try:
                        loop = asyncio.get_event_loop()
                        if loop.is_running():
                            # We can't await easily in a sync function without a bit of work if loop is running
                            import nest_asyncio
                            nest_asyncio.apply()
                        desc = asyncio.run(get_desc())
                        text += f"Scene Description:\n{desc}\n"
                    except:
                        pass
        except Exception as e:
            text += f"[Vision Error: {str(e)}]"
            
        return text

    def _parse_pdf(self, file_path: str) -> str:
        text = ""
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
                    else:
                        # Fallback to OCR for scanned pages
                        img = page.to_image()
                        text += pytesseract.image_to_string(img.original) + "\n"
        except Exception as e:
            text += f"[PDF Error: {str(e)}]"
        return text

    def _parse_docx(self, file_path: str) -> str:
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])

    def _parse_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _parse_excel(self, file_path: str) -> str:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = ""
        for sheet in wb.worksheets:
            text += f"--- Sheet: {sheet.title} ---\n"
            for row in sheet.iter_rows(values_only=True):
                row_str = " | ".join([str(c) if c is not None else "" for c in row])
                text += row_str + "\n"
        return text

    def _parse_delimited(self, file_path: str, delimiter: str) -> str:
        import csv
        rows = []
        with open(file_path, 'r', encoding='utf-8-sig', errors='replace', newline='') as source:
            for row in csv.reader(source, delimiter=delimiter):
                rows.append(' | '.join(cell.strip() for cell in row))
        return '\n'.join(rows)
