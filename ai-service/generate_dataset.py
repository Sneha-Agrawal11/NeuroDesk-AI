import os
import random
import string
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from reportlab.pdfgen import canvas

DATASET_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', 'benchmark_data_large'))
if not os.path.exists(DATASET_DIR):
    os.makedirs(DATASET_DIR)

def generate_random_text(words=100):
    return " ".join("".join(random.choices(string.ascii_lowercase, k=random.randint(3, 8))) for _ in range(words))

def create_pdf(path, text):
    c = canvas.Canvas(path)
    textobject = c.beginText(50, 800)
    for line in [text[i:i+80] for i in range(0, len(text), 80)]:
        textobject.textLine(line)
    c.drawText(textobject)
    c.save()

def create_docx(path, text):
    doc = Document()
    doc.add_paragraph(text)
    doc.save(path)

def create_pptx(path, text):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Presentation"
    slide.placeholders[1].text = text[:100]
    prs.save(path)

def create_xlsx(path, text):
    wb = Workbook()
    ws = wb.active
    words = text.split()
    for i in range(10):
        if i*3+2 < len(words):
            ws.append([words[i*3], words[i*3+1], words[i*3+2]])
    wb.save(path)

def create_image(path, text):
    img = Image.new('RGB', (400, 200), color=(255, 255, 255))
    d = ImageDraw.Draw(img)
    # Just draw the first few words to avoid font issues and fit
    d.text((10,10), text[:50], fill=(0,0,0))
    img.save(path)

def create_text(path, text):
    with open(path, 'w', encoding='utf-8') as f:
        f.write(text)

categories = ['resume', 'research_paper', 'invoice', 'notes', 'source_code']

# Keep track of texts to create duplicates
generated_texts = []

# Generate 500 files
for i in range(500):
    cat = random.choice(categories)
    
    # 10% chance of being an exact or near duplicate of a previous file
    if generated_texts and random.random() < 0.1:
        base_text = random.choice(generated_texts)
        if random.random() < 0.5:
            # Exact
            text = base_text
        else:
            # Near duplicate
            text = base_text + " " + generate_random_text(10)
    else:
        text = f"This is a {cat} document. " + generate_random_text(150)
        
    generated_texts.append(text)
    
    if i < 70:
        create_pdf(os.path.join(DATASET_DIR, f"doc_{i}.pdf"), text)
    elif i < 140:
        create_docx(os.path.join(DATASET_DIR, f"doc_{i}.docx"), text)
    elif i < 210:
        create_pptx(os.path.join(DATASET_DIR, f"doc_{i}.pptx"), text)
    elif i < 280:
        create_xlsx(os.path.join(DATASET_DIR, f"doc_{i}.xlsx"), text)
    elif i < 350:
        create_image(os.path.join(DATASET_DIR, f"doc_{i}.png"), text)
    elif i < 420:
        ext = random.choice(['.ts', '.js', '.jsx', '.tsx', '.json', '.csv'])
        create_text(os.path.join(DATASET_DIR, f"code_{i}{ext}"), text)
    else:
        ext = random.choice(['.txt', '.md'])
        create_text(os.path.join(DATASET_DIR, f"note_{i}{ext}"), text)

print(f"Generated 500 files in {DATASET_DIR}")
